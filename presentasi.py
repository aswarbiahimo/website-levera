from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_aesthetic_ppt():
    prs = Presentation()
    
    # --- KONFIGURASI WARNA (AESTHETIC PALETTE) ---
    # Background: Dark Grey (Elegant)
    BG_COLOR = RGBColor(30, 30, 30) 
    # Text Utama: Off-White
    TEXT_MAIN = RGBColor(240, 240, 240)
    # Accent: Coffee Gold (Mewah)
    ACCENT_COLOR = RGBColor(212, 175, 55)
    # Text Secondary: Grey
    TEXT_SEC = RGBColor(170, 170, 170)

    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_title_slide(prs, title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        set_background(slide)
        
        # Garis Aksen Vertikal
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(0.1), Inches(3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_COLOR
        shape.line.fill.background() # No border

        # Judul Besar
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(8), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Calibri Light'
        p.font.size = Pt(60)
        p.font.color.rgb = TEXT_MAIN
        p.font.bold = True

        # Subjudul
        txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(8), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.name = 'Calibri'
        p2.font.size = Pt(20)
        p2.font.color.rgb = ACCENT_COLOR

    def add_content_slide(prs, title_text, content_list):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)

        # Judul Slide (Kecil di pojok kiri atas - Minimalis)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.name = 'Calibri'
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_COLOR
        p.font.bold = True
        
        # Garis pembatas tipis
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(9), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(80, 80, 80)
        line.line.fill.background()

        # Konten
        top = 2.0
        for item in content_list:
            # Bullet point custom (Kotak kecil)
            bullet = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(top + 0.1), Inches(0.1), Inches(0.1))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = ACCENT_COLOR
            bullet.line.fill.background()

            # Teks
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(8.5), Inches(0.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item
            p.font.name = 'Calibri Light'
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_MAIN
            
            top += 0.8

    def add_card_slide(prs, title_text, cards_data):
        """Slide dengan kotak-kotak (Cards) agar terlihat mahal"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide)
        
        # Judul
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.size = Pt(24)
        p.font.color.rgb = ACCENT_COLOR
        p.font.bold = True

        left = 0.5
        for title, desc in cards_data:
            # Background Card
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.5), Inches(2.8), Inches(3.5))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(50, 50, 50) # Lighter grey for card
            shape.line.color.rgb = RGBColor(80, 80, 80)
            
            # Judul Card
            txBox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(2.7), Inches(2.4), Inches(0.8))
            p = txBox.text_frame.paragraphs[0]
            p.text = title
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.color.rgb = ACCENT_COLOR
            p.alignment = PP_ALIGN.CENTER
            
            # Deskripsi Card
            txBox2 = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.5), Inches(2.4), Inches(2.0))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            p2.text = desc
            p2.font.size = Pt(12)
            p2.font.color.rgb = TEXT_MAIN
            p2.alignment = PP_ALIGN.CENTER
            
            left += 3.1

    # --- MEMBUAT ISI SLIDE ---

    # 1. Title
    add_title_slide(prs, "LEVERA", "Modern Cafe & Resto Management System\nFinal Project Web Development")

    # 2. Introduction
    add_content_slide(prs, "Tentang Project", [
        "LEVERA adalah representasi digital dari pengalaman cafe modern.",
        "Menggabungkan estetika visual dengan fungsionalitas teknologi.",
        "Tujuan: Menciptakan ekosistem pemesanan yang seamless (tanpa hambatan) dari pelanggan ke dapur."
    ])

    # 3. Problem & Solution (Card Style)
    cards = [
        ("THE CHAOS", "Sistem manual menyebabkan antrean panjang, pesanan tertukar, dan selisih stok bahan baku."),
        ("THE GAP", "Pelanggan milenial menginginkan akses menu & reservasi instan tanpa harus menelepon."),
        ("THE SOLUTION", "Website terintegrasi untuk Reservasi Meja, E-Menu, dan Manajemen Order Real-time.")
    ]
    add_card_slide(prs, "Latar Belakang & Solusi", cards)

    # 4. Tech Stack (Card Style)
    tech_stack = [
        ("FRONTEND", "HTML5, CSS3, JavaScript.\nFramework: Bootstrap / Tailwind.\nFokus pada UI/UX Responsif."),
        ("BACKEND", "Bahasa: PHP / Python.\nLogika pemrosesan data pesanan dan manajemen user session."),
        ("DATABASE", "MySQL / MariaDB.\nMenyimpan data Menu, Transaksi, dan User secara terstruktur.")
    ]
    add_card_slide(prs, "Arsitektur Teknologi", tech_stack)

    # 5. Fitur Utama
    add_content_slide(prs, "Fitur Unggulan (Core Features)", [
        "Interactive Digital Menu: Tampilan menu dengan foto HD dan filter kategori.",
        "Smart Table Booking: Sistem pengecekan ketersediaan meja berdasarkan tanggal & waktu.",
        "Admin Dashboard: Panel kontrol pusat untuk memantau pesanan masuk & update stok.",
        "Responsive Layout: Akses nyaman melalui Smartphone (Mobile First) maupun Desktop."
    ])

    # 6. Placeholder untuk Demo
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    # Teks Tengah
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "LIVE DEMONSTRATION"
    p.font.size = Pt(44)
    p.font.color.rgb = ACCENT_COLOR
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # Subteks
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = "Silakan beralih ke browser untuk demo aplikasi"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_SEC
    p2.alignment = PP_ALIGN.CENTER

    # 7. Penutup
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Terima Kasih"
    p.font.size = Pt(50)
    p.font.color.rgb = ACCENT_COLOR
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Simpan
    filename = "Presentasi_LEVERA_Aesthetic.pptx"
    prs.save(filename)
    print(f"Berhasil! File '{filename}' telah dibuat dengan tema Dark Minimalist.")

if __name__ == "__main__":
    create_aesthetic_ppt()